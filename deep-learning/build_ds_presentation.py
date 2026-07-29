"""Build a DS1/DS2 experiment presentation from the repository template."""

from __future__ import annotations

import base64
import csv
import json
import math
import os
import sys
from dataclasses import dataclass
from pathlib import Path
from tempfile import TemporaryDirectory

import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches

from exp.analyze import Curve, RunRef, save_figure
from exp.plot_theme import apply_plot_theme


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "In-progress report.pptx"
OUTPUT = ROOT / "ds1_ds2_experiment_report.pptx"
TABLE_STYLE = "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"
STD_FINALS: dict[str, dict[str, tuple[int, float, float]]] = {}


@dataclass(frozen=True)
class Experiment:
    section: str
    title: str
    model_settings: tuple[tuple[str, str], ...]
    run_settings: tuple[tuple[str, str], ...]
    graph: Path
    result_headers: tuple[str, ...]
    result_rows: tuple[tuple[str, ...], ...]


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8", newline="") as file:
        return list(csv.DictReader(file))


def number(value: str | float, *, percent: bool = False, digits: int = 3) -> str:
    numeric = float(value)
    if percent:
        numeric *= 100
    if not math.isfinite(numeric):
        return "—"
    if abs(numeric) >= 1000:
        return f"{numeric:,.1f}"
    if percent:
        return f"{numeric:.2f}%"
    if abs(numeric) < 0.001 and numeric != 0:
        return f"{numeric:.2e}"
    return f"{numeric:.{digits}f}"


def mean_sd(
    title: str,
    series: str,
    *,
    percent: bool = False,
    digits: int = 3,
) -> str:
    run_count, mean, standard_deviation = STD_FINALS[title][series]
    if run_count == 10:
        return (
            f"{number(mean, percent=percent, digits=digits)} ± "
            f"{number(standard_deviation, percent=percent, digits=digits)}"
        )
    return f"{number(mean, percent=percent, digits=digits)} (n={run_count})"


def _row_count(path: Path) -> int:
    if not path.is_file():
        return 0
    with path.open(encoding="utf-8", errors="replace") as file:
        return max(0, sum(1 for _ in file) - 1)


def local_completed_seed_runs(
    _client,
    *,
    experiment_name: str,
    group_id: str,
    atomic_run_ids,
) -> dict[str, list[RunRef]]:
    """Select the richest successful local attempt for each condition and seed."""
    wanted = tuple(atomic_run_ids)
    grouped: dict[str, list[RunRef]] = {atomic: [] for atomic in wanted}
    root = ROOT / "exp" / experiment_name / "results" / "mlflow_artifacts"
    selected: dict[tuple[str, str], tuple[tuple[int, ...], RunRef]] = {}
    for condition_path in root.glob("*/config/condition.json"):
        artifact_root = condition_path.parent.parent
        try:
            condition = json.loads(condition_path.read_text(encoding="utf-8"))
            seed_data = json.loads((artifact_root / "config/seed.json").read_text(encoding="utf-8"))
        except (FileNotFoundError, json.JSONDecodeError, KeyError):
            continue
        atomic = str(condition.get("atomic_run_id", ""))
        if atomic not in grouped or condition.get("execution_group_id") != group_id:
            continue
        final_path = artifact_root / "metrics/final.json"
        if final_path.is_file():
            try:
                final = json.loads(final_path.read_text(encoding="utf-8"))
            except json.JSONDecodeError:
                continue
            if float(final.get("final/status/success", 1.0)) != 1.0:
                continue
        seed = str(seed_data["master"])
        score = (
            _row_count(artifact_root / "observations/source_curves.csv"),
            _row_count(artifact_root / "evaluations.csv"),
            _row_count(artifact_root / "updates.csv"),
            _row_count(artifact_root / "observations/activation_histogram.csv"),
            _row_count(artifact_root / "observations/attention.csv"),
        )
        run = RunRef(artifact_root.name, atomic, seed, 0, artifact_root)
        key = (atomic, seed)
        if key not in selected or score > selected[key][0]:
            selected[key] = (score, run)
    for (atomic, _seed), (_score, run) in selected.items():
        grouped[atomic].append(run)
    for runs_for_condition in grouped.values():
        runs_for_condition.sort(key=lambda run: int(run.seed))
    return grouped


def stored_completed_seed_runs(
    client,
    *,
    experiment_name: str,
    group_id: str,
    atomic_run_ids,
) -> dict[str, list[RunRef]]:
    """Resolve completed runs directly against this repository's artifact store."""
    wanted = tuple(atomic_run_ids)
    grouped: dict[str, list[RunRef]] = {atomic: [] for atomic in wanted}
    experiment = client.get_experiment_by_name(experiment_name)
    if experiment is None:
        return grouped
    runs = client.search_runs(
        [experiment.experiment_id],
        filter_string=f"attributes.status = 'FINISHED' and tags.`execution_group.id` = '{group_id}'",
        order_by=["attributes.start_time DESC"],
        max_results=5_000,
    )
    selected: dict[tuple[str, str], RunRef] = {}
    for run in runs:
        tags = run.data.tags
        atomic = tags.get("atomic_run.id", "")
        if tags.get("run.type") != "seed_trial" or atomic not in grouped:
            continue
        seed = str(run.data.params.get("seed/master", run.info.run_id))
        artifact_root = (
            ROOT
            / "infra/mlflow/data/artifacts"
            / str(experiment.experiment_id)
            / run.info.run_id
            / "artifacts"
        )
        selected.setdefault(
            (atomic, seed),
            RunRef(run.info.run_id, atomic, seed, int(run.info.start_time or 0), artifact_root),
        )
    for (atomic, _seed), run in selected.items():
        grouped[atomic].append(run)
    for runs_for_condition in grouped.values():
        runs_for_condition.sort(key=lambda run: int(run.seed))
    return grouped


def aggregate_sd(histories) -> Curve:
    """Aggregate common coordinates as mean ± sample standard deviation."""
    if not histories:
        return Curve.empty()
    common_steps = set(histories[0])
    for history in histories[1:]:
        common_steps.intersection_update(history)
    if not common_steps:
        return Curve.empty()
    steps = np.asarray(sorted(common_steps), dtype=float)
    values = np.asarray([[history[step] for step in steps] for history in histories], dtype=float)
    mean = values.mean(axis=0)
    standard_deviation = values.std(axis=0, ddof=1) if len(histories) > 1 else np.zeros_like(mean)
    return Curve(steps, mean, mean - standard_deviation, mean + standard_deviation, len(histories))


def make_sd_current_graphs(asset_dir: Path) -> dict[str, Path]:
    """Re-render repeated-run figures with standard-deviation error bands."""
    import exp.ds1.analyze.common as ds1_common
    import exp.ds2.analyze.common as ds2_common
    from exp.ds1.analyze.render import RENDERERS as DS1_RENDERERS
    from exp.ds2.analyze.render import RENDERERS as DS2_RENDERERS

    original = (
        ds1_common.aggregate,
        ds1_common.completed_seed_runs,
        ds2_common.aggregate,
        ds2_common.completed_seed_runs,
    )
    ds1_common.aggregate = aggregate_sd
    ds1_common.completed_seed_runs = local_completed_seed_runs
    ds2_common.aggregate = aggregate_sd
    ds2_common.completed_seed_runs = local_completed_seed_runs
    definitions = (
        ("ds1", "e01", "MNIST Optimizer Comparison"),
        ("ds1", "e02", "MNIST Weight Initialization"),
        ("ds1", "e03", "Weight Decay and Overfitting"),
        ("ds1", "e04", "Dropout and Overfitting"),
        ("ds1", "e05", "Batch Normalization and Weight Scale"),
        ("ds1", "e06", "Simple Convolutional Network"),
        ("ds1", "e06", "Deep Convolutional Network"),
        ("ds2", "e01", "Toy Word2Vec: CBOW vs Skip-gram"),
        ("ds2", "e03", "Small-Corpus RNN Language Model"),
        ("ds2", "e04", "Penn Treebank LSTM Language Model"),
    )
    outputs: dict[str, Path] = {}
    rendered: dict[tuple[str, str], tuple[Path, dict[str, Curve]]] = {}
    try:
        for domain, experiment_id, title in definitions:
            key = (domain, experiment_id)
            if key not in rendered:
                path = asset_dir / f"{domain}_{experiment_id}_mean_sd.png"
                renderer = DS1_RENDERERS[experiment_id] if domain == "ds1" else DS2_RENDERERS[experiment_id]
                figure, curves = renderer(None, "band", path)
                save_figure(figure, path)
                plt.close(figure)
                rendered[key] = (path, curves)
            path, curves = rendered[key]
            outputs[title] = path
            STD_FINALS[title] = {
                series: (
                    curve.run_count,
                    float(curve.mean[-1]),
                    float(curve.maximum[-1] - curve.mean[-1]),
                )
                for series, curve in curves.items()
                if len(curve.steps)
            }
        # This experiment stores its curve only in MLflow metric history.
        from exp.analyze import mlflow_client

        ds1_common.completed_seed_runs = original[1]
        spatial_path = asset_dir / "ds1_e08_mean_sd.png"
        tracking_uri = f"sqlite:///{(ROOT / 'infra/mlflow/data/mlflow.db').resolve()}"
        figure, curves = DS1_RENDERERS["e08"](mlflow_client(tracking_uri), "band", spatial_path)
        save_figure(figure, spatial_path)
        plt.close(figure)
        outputs["Spatial Layout Sensitivity"] = spatial_path
        STD_FINALS["Spatial Layout Sensitivity"] = {
            series: (
                curve.run_count,
                float(curve.mean[-1]),
                float(curve.maximum[-1] - curve.mean[-1]),
            )
            for series, curve in curves.items()
            if len(curve.steps)
        }
        # Completed Seq2Seq histories live in the MLflow artifact store rather
        # than the schema-v1 local mirror.
        ds2_common.completed_seed_runs = stored_completed_seed_runs
        client = mlflow_client(tracking_uri)
        for experiment_id, title in (
            ("e06", "Addition Seq2Seq Models"),
            ("e07", "Date Conversion Seq2Seq Models"),
        ):
            path = asset_dir / f"ds2_{experiment_id}_mean_sd.png"
            figure, curves = DS2_RENDERERS[experiment_id](client, "band", path)
            save_figure(figure, path)
            plt.close(figure)
            outputs[title] = path
            STD_FINALS[title] = {
                series: (
                    curve.run_count,
                    float(curve.mean[-1]),
                    float(curve.maximum[-1] - curve.mean[-1]),
                )
                for series, curve in curves.items()
                if len(curve.steps)
            }
    finally:
        (
            ds1_common.aggregate,
            ds1_common.completed_seed_runs,
            ds2_common.aggregate,
            ds2_common.completed_seed_runs,
        ) = original
    return outputs


def band_rows(
    path: Path,
    labels: dict[str, str],
    *,
    percent: bool = False,
) -> tuple[tuple[str, ...], ...]:
    output = []
    for row in read_csv(path):
        series = row["series"]
        if series not in labels:
            continue
        output.append(
            (
                labels[series],
                row["seed_runs"],
                number(row["final_mean"], percent=percent),
                f"{number(row['final_min'], percent=percent)}–{number(row['final_max'], percent=percent)}",
            )
        )
    return tuple(output)


def summary_rows(
    path: Path,
    labels: dict[str, str],
    metric_name: str,
    *,
    percent: bool = False,
) -> tuple[tuple[str, ...], ...]:
    rows = read_csv(path)
    grouped: dict[str, dict[str, dict[str, str]]] = {}
    for row in rows:
        grouped.setdefault(row["series"], {})[row["metric"]] = row
    output = []
    for series, label in labels.items():
        metrics = grouped.get(series, {})
        metric = metrics.get(metric_name)
        runtime = metrics.get("training_time_s")
        if metric is None:
            continue
        mean = number(metric["mean"], percent=False if metric.get("unit") == "percent" else percent)
        std = number(metric["standard_deviation"], percent=False if metric.get("unit") == "percent" else percent)
        if metric.get("unit") == "percent":
            mean += "%"
            std += " pp"
        runtime_text = "—" if runtime is None else f"{number(runtime['mean'], digits=1)} s"
        result = f"{mean} ± {std}" if metric["seed_runs"] == "10" else f"{mean} (n={metric['seed_runs']})"
        output.append((label, metric["seed_runs"], result, runtime_text))
    return tuple(output)


def pivot_train_test(
    path: Path,
    conditions: tuple[tuple[str, str], ...],
) -> tuple[tuple[str, ...], ...]:
    values = {row["series"]: row for row in read_csv(path)}
    output = []
    for prefix, label in conditions:
        train = values.get(f"{prefix}/train")
        test = values.get(f"{prefix}/test")
        output.append(
            (
                label,
                "—" if train is None else number(train["final_mean"], percent=True),
                "—" if test is None else number(test["final_mean"], percent=True),
                "—" if test is None else f"{number(test['final_min'], percent=True)}–{number(test['final_max'], percent=True)}",
            )
        )
    return tuple(output)


def batchnorm_rows(path: Path) -> tuple[tuple[str, ...], ...]:
    values = {row["series"]: row for row in read_csv(path)}
    standard_deviations = (
        1.0,
        0.541169526,
        0.292864456,
        0.158489319,
        0.085769590,
        0.046415888,
        0.025118864,
        0.013593564,
        0.007356423,
        0.003981072,
        0.002154435,
        0.001165914,
        0.000630957,
        0.000341455,
        0.000184785,
        0.0001,
    )
    rows = []
    for index, std in enumerate(standard_deviations, 1):
        off = values[f"BN-SCALE-{index:02d}-OFF"]
        on = values[f"BN-SCALE-{index:02d}-ON"]
        rows.append(
            (
                f"{std:.6g}",
                number(off["final_mean"], percent=True),
                number(on["final_mean"], percent=True),
            )
        )
    return tuple(rows)


def trajectory_rows(path: Path) -> tuple[tuple[str, ...], ...]:
    labels = {
        "TOY-SGD": "SGD",
        "TOY-MOMENTUM": "Momentum",
        "TOY-ADAGRAD": "AdaGrad",
        "TOY-ADAM": "Adam",
    }
    values = {row["series"]: float(row["final_mean"]) for row in read_csv(path)}
    rows = []
    for prefix, label in labels.items():
        x = values[f"{prefix}/x"]
        y = values[f"{prefix}/y"]
        objective = x * x / 20.0 + y * y
        rows.append((label, number(x), number(y), number(objective)))
    return tuple(rows)


def activation_rows(path: Path) -> tuple[tuple[str, ...], ...]:
    values = {row["series"]: row for row in read_csv(path)}
    labels = (
        ("SIGMOID", "STD001", "Sigmoid / σ=0.01"),
        ("SIGMOID", "XAVIER", "Sigmoid / Xavier"),
        ("SIGMOID", "HE", "Sigmoid / He"),
        ("SIGMOID", "STD1", "Sigmoid / σ=1.0"),
        ("TANH", "STD001", "Tanh / σ=0.01"),
        ("TANH", "XAVIER", "Tanh / Xavier"),
        ("TANH", "HE", "Tanh / He"),
        ("TANH", "STD1", "Tanh / σ=1.0"),
        ("RELU", "STD001", "ReLU / σ=0.01"),
        ("RELU", "XAVIER", "ReLU / Xavier"),
        ("RELU", "HE", "ReLU / He"),
        ("RELU", "STD1", "ReLU / σ=1.0"),
    )
    rows = []
    for activation, initializer, label in labels:
        prefix = f"ACT-{activation}-{initializer}"
        peaks = tuple(number(values[f"{prefix}/layer-{layer}"]["final_mean"], digits=0) for layer in range(1, 6))
        rows.append((label, *peaks))
    return tuple(rows)


def filter_rows(path: Path) -> tuple[tuple[str, ...], ...]:
    labels = {
        "CNN-SIMPLE-BOOK": "SimpleCNN (full training)",
        "CNN-SIMPLE-SPATIAL": "SimpleCNN (spatial)",
        "CNN-SIMPLE-SPATIAL-PERMUTED": "SimpleCNN (permuted)",
    }
    rows = []
    for row in read_csv(path):
        rows.append(
            (
                labels[row["condition"]],
                row["shape"],
                number(row["weight_min"]),
                number(row["weight_max"]),
                number(row["weight_std"]),
            )
        )
    return tuple(rows)


def make_word2vec_graph(path: Path) -> Path:
    rows = [
        row
        for row in read_csv(ROOT / "exp/ds2/results/image/e02_summary.csv")
        if row["metric"] == "final_loss"
    ]
    labels = {
        "W2V-PTB-CBOW-NS": "CBOW\nNeg. sampling",
        "W2V-PTB-SKIPGRAM-NS": "Skip-gram\nNeg. sampling",
        "W2V-PTB-CBOW-FULL": "CBOW\nFull softmax",
        "W2V-PTB-SKIPGRAM-FULL": "Skip-gram\nFull softmax",
    }
    figure, axis = plt.subplots(figsize=(10, 5.4))
    x = range(len(rows))
    means = [float(row["mean"]) for row in rows]
    deviations = [float(row["standard_deviation"]) for row in rows]
    axis.bar(x, means, yerr=deviations, capsize=4)
    axis.set_xticks(list(x), [labels[row["series"]] for row in rows])
    axis.set_ylabel("final loss")
    axis.set_title("PTB Word2Vec final loss")
    axis.grid(axis="y", alpha=0.25)
    figure.tight_layout()
    figure.savefig(path, dpi=180, transparent=True)
    plt.close(figure)
    return path


def make_activation_montage(path: Path) -> Path:
    source = ROOT / "exp/ds1/results/image/e10"
    files = [
        source / f"e10_band_act-{activation}-{initializer}.png"
        for activation in ("sigmoid", "tanh", "relu")
        for initializer in ("std001", "xavier", "he", "std1")
    ]
    opened = [Image.open(file).convert("RGB") for file in files]
    tile_width, tile_height = 640, 360
    canvas = Image.new("RGB", (tile_width * 4, tile_height * 3), "white")
    for index, image in enumerate(opened):
        fitted = ImageOps.contain(image, (tile_width, tile_height))
        x = (index % 4) * tile_width + (tile_width - fitted.width) // 2
        y = (index // 4) * tile_height + (tile_height - fitted.height) // 2
        canvas.paste(fitted, (x, y))
    canvas.save(path)
    for image in opened:
        image.close()
    return path


def make_message_figure(path: Path, heading: str, message: str) -> Path:
    figure, axis = plt.subplots(figsize=(7, 4.5))
    axis.set_title(heading)
    axis.text(0.5, 0.5, message, ha="center", va="center", wrap=True, transform=axis.transAxes)
    axis.set_xticks(())
    axis.set_yticks(())
    for spine in axis.spines.values():
        spine.set_visible(False)
    figure.tight_layout()
    figure.savefig(path, dpi=180)
    plt.close(figure)
    return path


def extract_notebook_figure(notebook: Path, image_index: int, output: Path) -> Path:
    document = json.loads(notebook.read_text(encoding="utf-8"))
    images = []
    for cell in document.get("cells", []):
        for item in cell.get("outputs", []):
            encoded = item.get("data", {}).get("image/png")
            if encoded:
                images.append("".join(encoded))
    output.write_bytes(base64.b64decode(images[image_index]))
    return output


def make_scalar_book_graph(path: Path, title: str, label: str, value: float, ylabel: str) -> Path:
    figure, axis = plt.subplots(figsize=(7, 4.5))
    axis.bar([label], [value])
    axis.set_title(title)
    axis.set_ylabel(ylabel)
    axis.grid(axis="y", alpha=0.25)
    figure.tight_layout()
    figure.savefig(path, dpi=180)
    plt.close(figure)
    return path


def make_ds1_cnn_book_graph(path: Path, deep: bool) -> Path:
    folder = (
        ROOT / "exp/ds1/results/original/data/e07/dlfs1.ch08.deep-convnet"
        if deep
        else ROOT / "exp/ds1/results/original/data/e06/dlfs1.ch07.simple-convnet"
    )
    rows = read_csv(folder / "metrics.csv")
    train = [(int(row["epoch"]), float(row["accuracy"])) for row in rows if row["split"] == "train"]
    test = [(int(row["epoch"]), float(row["accuracy"])) for row in rows if row["split"] == "test"]
    figure, axis = plt.subplots(figsize=(7, 4.5))
    axis.plot([x for x, _ in train], [y for _, y in train], marker="o", label="train")
    axis.plot([x for x, _ in test], [y for _, y in test], marker="s", label="test")
    axis.set_title("Book-code CNN accuracy")
    axis.set_xlabel("epoch")
    axis.set_ylabel("accuracy")
    axis.set_ylim(0.0, 1.02)
    axis.legend()
    figure.tight_layout()
    figure.savefig(path, dpi=180)
    plt.close(figure)
    return path


def make_original_attention(path: Path) -> Path:
    book = ROOT / "01_deep-learning-from-base/deep-learning-from-scratch-2"
    previous_cwd = Path.cwd()
    previous_path = list(sys.path)
    had_numpy_int = hasattr(np, "int")
    try:
        if not had_numpy_int:
            np.int = int  # type: ignore[attr-defined]
        os.chdir(book / "ch08")
        sys.path.insert(0, str(book / "ch08"))
        sys.path.insert(0, str(book))
        from dataset import sequence  # type: ignore[import-not-found]
        from attention_seq2seq import AttentionSeq2seq  # type: ignore[import-not-found]

        (_, _), (x_test, t_test) = sequence.load_data("date.txt")
        x_test = x_test[:, ::-1]
        char_to_id, id_to_char = sequence.get_vocab()
        model = AttentionSeq2seq(len(char_to_id), 16, 256)
        model.load_params(str(book / "ch08/AttentionSeq2seq.pkl"))
        np.random.seed(1984)
        index = int(np.random.randint(0, len(x_test)))
        x = x_test[[index]]
        t = t_test[[index]]
        model.forward(x, t)
        weights = np.asarray(model.decoder.attention.attention_weights)
        attention = weights.reshape(weights.shape[0], weights.shape[2])[:, ::-1]
        source = x[:, ::-1]
        row_labels = [id_to_char[int(item)] for item in source[0]]
        column_labels = [id_to_char[int(item)] for item in t[0][1:]]
    finally:
        if not had_numpy_int:
            delattr(np, "int")
        os.chdir(previous_cwd)
        sys.path[:] = previous_path
        for module in tuple(sys.modules):
            if module == "dataset" or module.startswith("dataset.") or module.startswith("common."):
                sys.modules.pop(module, None)
    figure, axis = plt.subplots(figsize=(7, 4.5))
    axis.pcolor(attention, cmap=plt.cm.Greys_r, vmin=0.0, vmax=1.0)
    axis.set_xticks(np.arange(attention.shape[1]) + 0.5, labels=row_labels)
    axis.set_yticks(np.arange(attention.shape[0]) + 0.5, labels=column_labels)
    axis.invert_yaxis()
    axis.set_title("Book-code attention alignment")
    figure.tight_layout()
    figure.savefig(path, dpi=180)
    plt.close(figure)
    return path


def make_comparison_graph(book_graph: Path, current_graph: Path, output: Path) -> Path:
    figure, axes = plt.subplots(1, 2, figsize=(13, 5.3))
    for axis, image_path, title in (
        (axes[0], book_graph, "Book original"),
        (axes[1], current_graph, "Current implementation (mean ± SD where n=10)"),
    ):
        with Image.open(image_path) as image:
            axis.imshow(image.convert("RGB"))
        axis.set_title(title)
        axis.set_xticks(())
        axis.set_yticks(())
        for spine in axis.spines.values():
            spine.set_visible(False)
    figure.tight_layout()
    figure.savefig(output, dpi=180)
    plt.close(figure)
    return output


def book_graphs(asset_dir: Path) -> dict[str, Path]:
    original = ROOT / "exp/ds1/results/original/image"
    notebooks = ROOT / "01_deep-learning-from-base/notebooks/b2"
    no_direct = make_message_figure(
        asset_dir / "no_direct_book_counterpart.png",
        "No direct book experiment",
        "This condition is a project extension.\nThe closest book component is identified in the setup and table.",
    )
    not_cached = make_message_figure(
        asset_dir / "book_result_not_cached.png",
        "Book condition exists",
        "The original source defines this condition,\nbut a completed original run is not cached in the repository.",
    )
    return {
        "MNIST Optimizer Comparison": original / "e01_optimizer_compare_mnist.png",
        "MNIST Weight Initialization": original / "e02_weight_init_compare.png",
        "Weight Decay and Overfitting": original / "e03_overfit_weight_decay.png",
        "Dropout and Overfitting": original / "e04_overfit_dropout.png",
        "Batch Normalization and Weight Scale": original / "e05_batch_norm_test.png",
        "Simple Convolutional Network": original / "e06_train_convnet.png",
        "Deep Convolutional Network": make_ds1_cnn_book_graph(asset_dir / "book_deep_cnn.png", True),
        "Spatial Layout Sensitivity": no_direct,
        "Optimizer Trajectories": original / "e09_optimizer_compare_naive.png",
        "Activation Distribution Observation": original / "e10_weight_init_activation_histogram.png",
        "First-Layer Convolution Filters": original / "e11_filters_trained.png",
        "Toy Word2Vec: CBOW vs Skip-gram": ROOT / "exp/ds2/results/original/image/e01_toy_cbow.png",
        "PTB Word2Vec Objectives": extract_notebook_figure(
            notebooks / "03_word2vec.ipynb", 2, asset_dir / "book_ptb_cbow.png"
        ),
        "Small-Corpus RNN Language Model": extract_notebook_figure(
            notebooks / "04_RNN.ipynb", 0, asset_dir / "book_small_rnn.png"
        ),
        "Penn Treebank LSTM Language Model": make_scalar_book_graph(
            asset_dir / "book_lstm_ppl.png", "Book checkpoint test perplexity", "LSTM", 136.078, "perplexity"
        ),
        "Penn Treebank LM Recipes": make_scalar_book_graph(
            asset_dir / "book_better_lstm_ppl.png", "Book-code improved LSTM", "Improved LSTM", 80.826, "test perplexity"
        ),
        "Addition Seq2Seq Models": extract_notebook_figure(
            notebooks / "06_seq2seq.ipynb", 1, asset_dir / "book_addition_seq2seq.png"
        ),
        "Date Conversion Seq2Seq Models": extract_notebook_figure(
            notebooks / "07_attention.ipynb", 1, asset_dir / "book_date_attention.png"
        ),
        "Attention Alignment Observation": make_original_attention(asset_dir / "book_attention_alignment.png"),
        "_not_cached": not_cached,
    }


def book_settings(title: str) -> tuple[tuple[str, str], ...]:
    settings = {
        "MNIST Optimizer Comparison": (
            ("Model", "MLP: 784 → 100 × 4 → 10"),
            ("Data / batch", "MNIST full train / 128"),
            ("Budget", "2,000 updates"),
            ("Variable", "SGD, Momentum, AdaGrad, Adam"),
            ("Learning rates", "0.01; Adam 0.001"),
            ("Result", "Raw mini-batch loss"),
        ),
        "MNIST Weight Initialization": (
            ("Model", "MLP: 784 → 100 × 4 → 10, ReLU"),
            ("Data / batch", "MNIST full train / 128"),
            ("Budget", "2,000 updates"),
            ("Variable", "σ=0.01, Xavier, He"),
            ("Optimizer", "SGD, learning rate 0.01"),
            ("Result", "Raw mini-batch loss"),
        ),
        "Weight Decay and Overfitting": (
            ("Model", "MLP: 784 → 100 × 6 → 10"),
            ("Data / batch", "First 300 train, full test / 100"),
            ("Budget", "201 epoch observations"),
            ("Book-default condition", "L2 weight decay 0.1"),
            ("Optimizer", "SGD, learning rate 0.01"),
            ("Result", "Train/test accuracy"),
        ),
        "Dropout and Overfitting": (
            ("Model", "MLP: 784 → 100 × 6 → 10"),
            ("Data / batch", "First 300 train, full test / 100"),
            ("Budget", "301 epochs"),
            ("Book-default condition", "Dropout 0.2"),
            ("Optimizer", "SGD, learning rate 0.01"),
            ("Result", "Train/test accuracy"),
        ),
        "Batch Normalization and Weight Scale": (
            ("Model", "MLP: 784 → 100 × 5 → 10"),
            ("Data / batch", "First 1,000 train / 100"),
            ("Budget", "20 epoch observations"),
            ("Variable", "BN off/on; 16 weight scales"),
            ("Optimizer", "SGD, learning rate 0.01"),
            ("Result", "Training accuracy"),
        ),
        "Simple Convolutional Network": (
            ("Model", "Conv30 5×5 → pool → FC100 → 10"),
            ("Data / batch", "MNIST images / 100"),
            ("Budget", "20 epochs"),
            ("Optimizer", "Adam, learning rate 0.001"),
            ("Evaluation", "First 1,000 per epoch; full test at end"),
            ("Seed", "Single book-code run"),
        ),
        "Deep Convolutional Network": (
            ("Model", "Six convolutions → FC50 → 10"),
            ("Data / batch", "MNIST images / 100"),
            ("Budget", "20 epochs"),
            ("Optimizer", "Adam, learning rate 0.001"),
            ("Evaluation", "First 1,000 per epoch; full test at end"),
            ("Seed", "Single book-code run"),
        ),
        "Spatial Layout Sensitivity": (
            ("Direct counterpart", "None"),
            ("Closest book components", "MLP and SimpleCNN experiments"),
            ("Project extension", "Fixed pixel permutation"),
            ("Book comparison scope", "Architecture baseline only"),
        ),
        "Optimizer Trajectories": (
            ("Objective / start", "x²/20 + y² / (-7, 2)"),
            ("Budget", "30 updates"),
            ("Optimizers", "SGD, Momentum, AdaGrad, Adam"),
            ("Learning rates", "0.95, 0.1, 1.5, 0.3"),
            ("Result", "Pre-update x/y trajectory"),
            ("Seed", "Deterministic analytic run"),
        ),
        "Activation Distribution Observation": (
            ("Model", "Five layers, width 100, bias 0"),
            ("Input", "1,000 × 100 standard normal"),
            ("Book-default condition", "Sigmoid, weight σ=1.0"),
            ("Training", "None; forward observation"),
            ("Result", "Five activation histograms"),
            ("Seed", "NumPy seed 1"),
        ),
        "First-Layer Convolution Filters": (
            ("Model", "Book SimpleCNN"),
            ("Observed parameter", "First 30 kernels, 1×5×5"),
            ("States", "Initial and trained"),
            ("Training", "20 epochs on MNIST"),
            ("Result", "Grayscale filter mosaics"),
        ),
        "Toy Word2Vec: CBOW vs Skip-gram": (
            ("Book model", "SimpleCBOW only"),
            ("Corpus / window", "Toy sentence / 1"),
            ("Embedding / batch", "5 / 3"),
            ("Budget", "1,000 epochs"),
            ("Objective / optimizer", "Full softmax / Adam 0.001"),
            ("Extension in current run", "Toy Skip-gram"),
        ),
        "PTB Word2Vec Objectives": (
            ("Book models", "CBOW and selectable Skip-gram"),
            ("Corpus / window", "PTB train / 5"),
            ("Embedding / batch", "100 / 100"),
            ("Budget", "10 epochs"),
            ("Book objective", "Negative sampling, 5 negatives"),
            ("Extension in current run", "Full-softmax variants"),
        ),
        "Small-Corpus RNN Language Model": (
            ("Model", "SimpleRnnlm, embedding/hidden 100"),
            ("Corpus", "First 1,000 PTB tokens"),
            ("Batch / BPTT", "10 / 5"),
            ("Budget", "100 epochs"),
            ("Optimizer", "SGD, learning rate 0.1"),
            ("Result", "Training perplexity"),
        ),
        "Penn Treebank LSTM Language Model": (
            ("Model", "LSTM Rnnlm, embedding/hidden 100"),
            ("Corpus", "PTB train/test"),
            ("Batch / BPTT", "20 / 35"),
            ("Budget", "4 epochs"),
            ("Optimizer / clip", "SGD 20 / 0.25"),
            ("Result", "Test perplexity"),
        ),
        "Penn Treebank LM Recipes": (
            ("Direct combined experiment", "None"),
            ("Closest book recipe", "Improved LSTM, 650 units"),
            ("Budget", "40 epochs"),
            ("Optimizer / clip", "SGD 20 / 0.25"),
            ("Schedule", "Validation-driven LR decay"),
            ("Current extension", "Three-recipe comparison"),
        ),
        "Addition Seq2Seq Models": (
            ("Models", "Vanilla / Peeky selectable"),
            ("Data", "Addition, 90/10 split"),
            ("Embedding / hidden", "16 / 128"),
            ("Batch / budget", "128 / 25 epochs"),
            ("Optimizer / clip", "Adam 0.001 / 5"),
            ("Variable", "Forward vs reversed input"),
        ),
        "Date Conversion Seq2Seq Models": (
            ("Models", "Vanilla, Peeky, Attention selectable"),
            ("Data / order", "Date conversion / reversed"),
            ("Embedding / hidden", "16 / 256"),
            ("Batch / budget", "128 / 10 epochs"),
            ("Optimizer / clip", "Adam 0.001 / 5"),
            ("Result", "Full-test exact match"),
        ),
        "Attention Alignment Observation": (
            ("Model", "Trained AttentionSeq2seq"),
            ("Data / order", "Date test / reversed"),
            ("Selection seed", "1984"),
            ("Examples", "5"),
            ("Result", "Decoder × encoder attention matrices"),
            ("Book checkpoint", "Available"),
        ),
    }
    return settings[title]


def reproduction_setup(title: str) -> tuple[tuple[tuple[str, str], ...], str]:
    setups = {
        "MNIST Optimizer Comparison": (
            (
                ("Dataset", "MNIST full training set; flattened 784-dimensional input; batch 128; sampling with replacement"),
                ("Model architecture", "MLP 784 → 100 → 100 → 100 → 100 → 10; ReLU; He initialization"),
                ("Optimizer & hyperparameters", "SGD 0.01; Momentum 0.01 with momentum 0.9; AdaGrad 0.01; Adam 0.001; mean softmax cross-entropy"),
                ("Updates / epochs", "2,000 optimizer updates, capped at 5 epochs; loss recorded per update; 10 paired seeds"),
            ),
            "Difference from book: the core architecture, optimizer settings, batch size, and 2,000-update budget match; the reproduction aggregates 10 paired seeds instead of one source run.",
        ),
        "MNIST Weight Initialization": (
            (
                ("Dataset", "MNIST full training set; flattened 784-dimensional input; batch 128; sampling with replacement"),
                ("Model architecture", "MLP 784 → 100 → 100 → 100 → 100 → 10; ReLU"),
                ("Optimizer & hyperparameters", "Initializers: normal σ=0.01, Xavier, and He; SGD 0.01; mean softmax cross-entropy"),
                ("Updates / epochs", "2,000 optimizer updates, capped at 5 epochs; loss recorded per update; 10 paired seeds"),
            ),
            "Difference from book: the three initialization conditions and training budget match; the reproduction reports a 10-seed aggregate rather than a single source run.",
        ),
        "Weight Decay and Overfitting": (
            (
                ("Dataset", "First 300 MNIST training examples and full 10,000-example test set; batch 100; sampling with replacement"),
                ("Model architecture", "MLP 784 → 100 × 6 → 10; ReLU; He initialization"),
                ("Optimizer & hyperparameters", "SGD 0.01; L2 weight decay 0.0 and 0.1; mean softmax cross-entropy"),
                ("Updates / epochs", "201 epochs, up to 601 updates; train/test accuracy every 3 updates; 10 paired seeds"),
            ),
            "Difference from book: weight decay 0.1 reproduces the source-default run; the no-decay control and 10-seed aggregation are added comparison conditions.",
        ),
        "Dropout and Overfitting": (
            (
                ("Dataset", "First 300 MNIST training examples and full 10,000-example test set; batch 100; sampling with replacement"),
                ("Model architecture", "MLP 784 → 100 × 6 → 10; ReLU; He initialization"),
                ("Optimizer & hyperparameters", "SGD 0.01; dropout ratio 0.0 and 0.2; mean softmax cross-entropy"),
                ("Updates / epochs", "301 epochs, up to 903 updates; train/test accuracy at the first update of each epoch; 10 paired seeds"),
            ),
            "Difference from book: dropout 0.2 reproduces the source-default run; the dropout-off control and 10-seed aggregation are added.",
        ),
        "Batch Normalization and Weight Scale": (
            (
                ("Dataset", "First 1,000 MNIST training examples; flattened input; batch 100; sampling with replacement"),
                ("Model architecture", "MLP 784 → 100 × 5 → 10; ReLU; BatchNorm enabled or disabled"),
                ("Optimizer & hyperparameters", "SGD 0.01; 16 normal weight scales logarithmically spaced from σ=1.0 to 0.0001"),
                ("Updates / epochs", "20 epochs; training accuracy every 10 updates through update 191; 10 paired seeds"),
            ),
            "Difference from book: model, weight-scale grid, BatchNorm axis, and cadence match; the reproduction repeats every condition across 10 paired seeds.",
        ),
        "Simple Convolutional Network": (
            (
                ("Dataset", "MNIST images, 60,000 train and 10,000 test; batch 100; sampling with replacement"),
                ("Model architecture", "Conv 1→30 with 5×5 kernels → pooling → fully connected 100 → 10; normal initialization σ=0.01"),
                ("Optimizer & hyperparameters", "Adam 0.001; mean softmax cross-entropy; CuPy GPU; float64"),
                ("Updates / epochs", "20 epochs; first-1,000 train/test evaluation each epoch; full test evaluation at termination; 10 seeds"),
            ),
            "Difference from book: architecture and learning protocol match; the reproduction uses multi-seed aggregation and the repository GPU/numerics backend.",
        ),
        "Deep Convolutional Network": (
            (
                ("Dataset", "MNIST images, 60,000 train and 10,000 test; batch 100; sampling with replacement"),
                ("Model architecture", "Six convolution layers with channel stages 16/32/64 → fully connected 50 → 10; dropout 0.5; He initialization"),
                ("Optimizer & hyperparameters", "Adam 0.001; mean softmax cross-entropy; CuPy GPU; float64"),
                ("Updates / epochs", "20 epochs; first-1,000 train/test evaluation each epoch; full test evaluation at termination; 10 seeds"),
            ),
            "Difference from book: architecture and learning protocol match; the reproduction uses multi-seed aggregation and the repository GPU/numerics backend.",
        ),
        "Spatial Layout Sensitivity": (
            (
                ("Dataset", "MNIST full train/test; original pixels or one fixed permutation with seed 20260808; batch 100"),
                ("Model architecture", "Parameter-matched MLP 784→489→100→10 versus SimpleCNN Conv30 5×5→FC100→10"),
                ("Optimizer & hyperparameters", "Adam 0.001; ReLU; model-specific He or normal σ=0.01 initialization; mean softmax cross-entropy"),
                ("Updates / epochs", "2 epochs; evaluation every 20 updates and full test at each epoch end; 10 paired seeds"),
            ),
            "Difference from book: this is a project extension with no direct book experiment; it combines book-style MLP/CNN components with a fixed spatial permutation.",
        ),
        "Optimizer Trajectories": (
            (
                ("Dataset", "No sampled dataset; analytic objective f(x,y)=x²/20+y²; initial point (-7, 2)"),
                ("Model architecture", "Two-parameter analytic quadratic observation model"),
                ("Optimizer & hyperparameters", "SGD 0.95; Momentum 0.1 with momentum 0.9; AdaGrad 1.5; Adam 0.3; float64"),
                ("Updates / epochs", "60 optimizer updates; x, y, objective, and gradient recorded at every step; deterministic single run"),
            ),
            "Difference from book: optimizer settings match, but the book records 30 pre-update positions while the reproduction extends the trajectory to 60 updates.",
        ),
        "Activation Distribution Observation": (
            (
                ("Dataset", "1,000 × 100 synthetic standard-normal inputs; fixed input seed 40402"),
                ("Model architecture", "Five-layer width-100 MLP observation network; zero bias; fixed model seed 40403"),
                ("Optimizer & hyperparameters", "No optimizer; activations sigmoid/tanh/ReLU crossed with σ=0.01, Xavier, He, and σ=1.0 initialization"),
                ("Updates / epochs", "No training; one forward observation; five per-layer activation histograms per condition"),
            ),
            "Difference from book: the book source-default figure is sigmoid with σ=1.0 and NumPy seed 1; the reproduction expands this to a 3×4 activation/initializer grid with explicit seeds.",
        ),
        "First-Layer Convolution Filters": (
            (
                ("Dataset", "MNIST context from the completed SimpleCNN and spatial-layout training runs"),
                ("Model architecture", "First convolution tensor from SimpleCNN: 30 output filters × 1 input channel × 5×5 kernel"),
                ("Optimizer & hyperparameters", "No additional optimization; final checkpoints are observed on one shared grayscale weight scale"),
                ("Updates / epochs", "Post-training observation at the final checkpoint; master seed 1; three current training/input conditions"),
            ),
            "Difference from book: the book contrasts initial and trained SimpleCNN filters; the reproduction compares trained filters from full, spatial, and spatial-permuted conditions.",
        ),
        "Toy Word2Vec: CBOW vs Skip-gram": (
            (
                ("Dataset", 'Toy corpus: "You say goodbye and I say hello."; context window 1; batch 3'),
                ("Model architecture", "Simple CBOW and Simple Skip-gram; embedding dimension 5"),
                ("Optimizer & hyperparameters", "Adam 0.001; full-softmax objective; float32 deterministic CuPy execution"),
                ("Updates / epochs", "1,000 epochs; mean loss recorded every 20 updates; 10 paired seeds"),
            ),
            "Difference from book: CBOW reproduces the book experiment; toy Skip-gram and the 10-seed aggregate are project extensions.",
        ),
        "PTB Word2Vec Objectives": (
            (
                ("Dataset", "Penn Treebank training corpus; context window 5; batch 100; drop-last batching"),
                ("Model architecture", "CBOW and Skip-gram; embedding dimension 100"),
                ("Optimizer & hyperparameters", "Adam 0.001; negative sampling with 5 negatives and unigram power 0.75, plus full-softmax variants"),
                ("Updates / epochs", "10 epochs; interval-mean loss every 20 updates; 10 paired deterministic GPU seeds"),
            ),
            "Difference from book: CBOW/Skip-gram negative-sampling conditions follow the book; both full-softmax conditions and multi-seed aggregation are extensions.",
        ),
        "Small-Corpus RNN Language Model": (
            (
                ("Dataset", "First 1,000 Penn Treebank training tokens; batch 10; truncated sequence length 5"),
                ("Model architecture", "Vanilla RNN language model; word-vector dimension 100; hidden dimension 100"),
                ("Optimizer & hyperparameters", "SGD 0.1; temporal softmax cross-entropy; no gradient clipping; float32 deterministic CuPy"),
                ("Updates / epochs", "100 epochs; training perplexity aggregated and recorded every 20 updates; 10 seeds"),
            ),
            "Difference from book: data, model, optimizer, and epoch budget match; the reproduction runs 10 deterministic GPU seeds instead of one source run.",
        ),
        "Penn Treebank LSTM Language Model": (
            (
                ("Dataset", "Penn Treebank train/test; batch 20; truncated sequence length 35"),
                ("Model architecture", "LSTM language model; word-vector dimension 100; hidden dimension 100"),
                ("Optimizer & hyperparameters", "SGD 20.0; gradient clipping 0.25; temporal softmax cross-entropy; float32 deterministic CuPy"),
                ("Updates / epochs", "4 epochs; train perplexity every 20 updates; full test perplexity at termination; 10 seeds"),
            ),
            "Difference from book: the model and training protocol match; the reproduction aggregates 10 deterministic GPU runs rather than one checkpoint run.",
        ),
        "Penn Treebank LM Recipes": (
            (
                ("Dataset", "Penn Treebank train/validation/test; batch 20; truncated sequence length 35"),
                ("Model architecture", "Vanilla RNN, LSTM, and improved LSTM; word-vector/hidden dimensions 650; improved model dropout 0.5"),
                ("Optimizer & hyperparameters", "SGD 20.0; gradient clipping 0.25; validation-triggered learning-rate decay by factor 4"),
                ("Updates / epochs", "Up to 40 epochs; train perplexity every 20 updates; validation each epoch; test at termination"),
            ),
            "Difference from book: this unified three-recipe comparison is a project extension assembled from book architectures; the book does not emit one combined comparison run.",
        ),
        "Addition Seq2Seq Models": (
            (
                ("Dataset", "Addition sequence dataset; legacy 90/10 split with seed 1984; forward or reversed inputs; batch 128"),
                ("Model architecture", "Vanilla Seq2Seq and Peeky Seq2Seq; embedding dimension 16; hidden dimension 128"),
                ("Optimizer & hyperparameters", "Adam 0.001; gradient clipping 5.0; temporal softmax cross-entropy; greedy decoding"),
                ("Updates / epochs", "25 epochs; full-test exact match every epoch; 10 fixed predictions; 10 paired seeds"),
            ),
            "Difference from book: the individual model/order conditions follow source-selectable book variants; the reproduction executes them as one paired 10-seed comparison.",
        ),
        "Date Conversion Seq2Seq Models": (
            (
                ("Dataset", "Date-format conversion dataset; legacy 90/10 split with seed 1984; reversed inputs; batch 128"),
                ("Model architecture", "Vanilla, Peeky, and Attention Seq2Seq; embedding dimension 16; hidden dimension 256"),
                ("Optimizer & hyperparameters", "Adam 0.001; gradient clipping 5.0; temporal softmax cross-entropy; greedy decoding"),
                ("Updates / epochs", "10 epochs; full-test exact match every epoch; 10 fixed predictions; 10 paired seeds"),
            ),
            "Difference from book: the three source-selectable architectures are combined into one comparison and repeated over the same 10 paired seeds; the book reports selected single runs.",
        ),
        "Attention Alignment Observation": (
            (
                ("Dataset", "Date conversion test split; reversed input; selection seed 1984; five selected examples"),
                ("Model architecture", "Attention Seq2Seq with embedding dimension 16 and hidden dimension 256"),
                ("Optimizer & hyperparameters", "No optimizer; greedy decoding from a completed trained attention checkpoint; float32 CPU observation"),
                ("Updates / epochs", "No training; one forward/decode pass per selected example; decoder×encoder attention matrix recorded"),
            ),
            "Difference from book: the observation protocol matches, but the current implementation has no completed observation run because its trained checkpoint dependency is unavailable.",
        ),
    }
    return setups[title]


def comparison_spec(title: str) -> tuple[str, str, str]:
    specs = {
        "MNIST Optimizer Comparison": (
            "Optimizer",
            "SGD vs Momentum vs AdaGrad vs Adam → training loss",
        ),
        "MNIST Weight Initialization": (
            "Weight initialization",
            "Normal σ=0.01 vs Xavier vs He → training loss",
        ),
        "Weight Decay and Overfitting": (
            "L2 weight decay",
            "0.0 vs 0.1 → train/test accuracy and generalization gap",
        ),
        "Dropout and Overfitting": (
            "Dropout",
            "Off (0.0) vs on (0.2) → train/test accuracy and generalization gap",
        ),
        "Batch Normalization and Weight Scale": (
            "BatchNorm × initial weight scale",
            "Off/on × 16 values of σ (1.0 to 0.0001) → training accuracy",
        ),
        "Simple Convolutional Network": (
            "Implementation / run protocol",
            "Book single run vs same-architecture 10-seed reproduction → full-test accuracy",
        ),
        "Deep Convolutional Network": (
            "Implementation / run protocol",
            "Book single run vs same-architecture 10-seed reproduction → full-test accuracy",
        ),
        "Spatial Layout Sensitivity": (
            "Model family × input layout",
            "MLP/CNN × original/fixed-permuted pixels (2×2) → train/test accuracy",
        ),
        "Optimizer Trajectories": (
            "Optimizer",
            "SGD vs Momentum vs AdaGrad vs Adam → parameter path and final objective",
        ),
        "Activation Distribution Observation": (
            "Activation × initialization",
            "Sigmoid/tanh/ReLU × σ=0.01/Xavier/He/σ=1.0 (3×4) → layer histograms",
        ),
        "First-Layer Convolution Filters": (
            "Training / input condition",
            "Full-data vs spatial-original vs spatial-permuted checkpoints → first-layer filters",
        ),
        "Toy Word2Vec: CBOW vs Skip-gram": (
            "Architecture",
            "Simple CBOW vs Simple Skip-gram → final full-softmax loss",
        ),
        "PTB Word2Vec Objectives": (
            "Architecture × objective",
            "CBOW/Skip-gram × negative sampling/full softmax (2×2) → final loss",
        ),
        "Small-Corpus RNN Language Model": (
            "Implementation / run protocol",
            "Book single run vs 10-seed reproduction → final training perplexity",
        ),
        "Penn Treebank LSTM Language Model": (
            "Implementation / run protocol",
            "Book checkpoint vs 10-seed reproduction → final test perplexity",
        ),
        "Penn Treebank LM Recipes": (
            "Language-model recipe",
            "Vanilla RNN vs LSTM vs improved LSTM → final test perplexity",
        ),
        "Addition Seq2Seq Models": (
            "Architecture × input order",
            "Vanilla/Peeky × forward/reversed input (2×2) → exact-match accuracy",
        ),
        "Date Conversion Seq2Seq Models": (
            "Architecture",
            "Vanilla vs Peeky vs Attention (reversed input fixed) → exact-match accuracy",
        ),
        "Attention Alignment Observation": (
            "Implementation / checkpoint availability",
            "Book trained checkpoint vs current observation dependency → attention alignment maps",
        ),
    }
    variable, levels_and_outcome = specs[title]
    levels, outcome = levels_and_outcome.rsplit(" → ", 1)
    return variable, levels, outcome


def separated_setup(
    title: str,
) -> tuple[tuple[str, str, str], tuple[tuple[str, str], ...], str]:
    rows, difference = reproduction_setup(title)
    replacements: dict[str, dict[str, tuple[str, str]]] = {
        "MNIST Optimizer Comparison": {
            "Optimizer & hyperparameters": (
                "Objective / numerics",
                "Mean softmax cross-entropy; NumPy CPU; float64",
            ),
        },
        "MNIST Weight Initialization": {
            "Optimizer & hyperparameters": (
                "Optimizer / objective",
                "SGD 0.01; mean softmax cross-entropy",
            ),
        },
        "Weight Decay and Overfitting": {
            "Optimizer & hyperparameters": (
                "Optimizer / objective",
                "SGD 0.01; mean softmax cross-entropy",
            ),
        },
        "Dropout and Overfitting": {
            "Optimizer & hyperparameters": (
                "Optimizer / objective",
                "SGD 0.01; mean softmax cross-entropy",
            ),
        },
        "Batch Normalization and Weight Scale": {
            "Model architecture": (
                "Common model structure",
                "MLP 784 → 100 × 5 → 10; ReLU",
            ),
            "Optimizer & hyperparameters": (
                "Optimizer",
                "SGD 0.01",
            ),
        },
        "Simple Convolutional Network": {
            "Optimizer & hyperparameters": (
                "Optimizer / objective",
                "Adam 0.001; mean softmax cross-entropy",
            ),
            "Updates / epochs": (
                "Training budget / evaluation",
                "20 epochs; first-1,000 train/test each epoch; full test at termination",
            ),
        },
        "Deep Convolutional Network": {
            "Optimizer & hyperparameters": (
                "Optimizer / objective",
                "Adam 0.001; mean softmax cross-entropy",
            ),
            "Updates / epochs": (
                "Training budget / evaluation",
                "20 epochs; first-1,000 train/test each epoch; full test at termination",
            ),
        },
        "Spatial Layout Sensitivity": {
            "Dataset": (
                "Dataset",
                "MNIST full train/test; one shared permutation seed 20260808; batch 100",
            ),
            "Model architecture": (
                "Common prediction task",
                "784-dimensional input to 10 classes; model family is the compared factor",
            ),
            "Optimizer & hyperparameters": (
                "Optimizer / objective",
                "Adam 0.001; mean softmax cross-entropy",
            ),
        },
        "Optimizer Trajectories": {
            "Optimizer & hyperparameters": (
                "Gradient / numerics",
                "Analytic gradient; NumPy CPU; float64",
            ),
        },
        "Activation Distribution Observation": {
            "Optimizer & hyperparameters": (
                "Training",
                "No optimizer; forward observation only",
            ),
        },
        "First-Layer Convolution Filters": {
            "Dataset": (
                "Dataset context",
                "MNIST; checkpoint condition is the compared factor",
            ),
            "Updates / epochs": (
                "Observation point",
                "Final post-training checkpoint; master seed 1",
            ),
        },
        "Toy Word2Vec: CBOW vs Skip-gram": {
            "Model architecture": (
                "Common representation size",
                "Embedding dimension 5",
            ),
        },
        "PTB Word2Vec Objectives": {
            "Model architecture": (
                "Common representation size",
                "Embedding dimension 100",
            ),
            "Optimizer & hyperparameters": (
                "Optimizer",
                "Adam 0.001; deterministic float32 CuPy execution",
            ),
        },
        "Small-Corpus RNN Language Model": {
            "Optimizer & hyperparameters": (
                "Optimizer / objective",
                "SGD 0.1; temporal softmax cross-entropy; no gradient clipping",
            ),
            "Updates / epochs": (
                "Training budget / evaluation",
                "100 epochs; training perplexity aggregated every 20 updates",
            ),
        },
        "Penn Treebank LSTM Language Model": {
            "Optimizer & hyperparameters": (
                "Optimizer / objective",
                "SGD 20.0; gradient clipping 0.25; temporal softmax cross-entropy",
            ),
            "Updates / epochs": (
                "Training budget / evaluation",
                "4 epochs; train perplexity every 20 updates; full test at termination",
            ),
        },
        "Penn Treebank LM Recipes": {
            "Model architecture": (
                "Common representation size",
                "Word-vector and hidden dimensions 650",
            ),
        },
        "Addition Seq2Seq Models": {
            "Dataset": (
                "Dataset",
                "Addition dataset; shared legacy 90/10 split seed 1984; batch 128",
            ),
            "Model architecture": (
                "Common dimensions",
                "Embedding 16; hidden 128",
            ),
        },
        "Date Conversion Seq2Seq Models": {
            "Model architecture": (
                "Common dimensions / input",
                "Embedding 16; hidden 256; reversed input",
            ),
        },
        "Attention Alignment Observation": {
            "Optimizer & hyperparameters": (
                "Decoding",
                "No optimizer; greedy decoding from a trained attention checkpoint",
            ),
        },
    }
    fixed_rows = tuple(
        replacements.get(title, {}).get(component, (component, protocol))
        for component, protocol in rows
    )
    return comparison_spec(title), fixed_rows, difference


def original_last_value(folder: Path, condition: str | None = None, *, split: str | None = None) -> float:
    rows = read_csv(folder / "metrics.csv")
    selected = rows
    if condition is not None:
        selected = [row for row in selected if row.get("condition") == condition]
    if split is not None:
        selected = [row for row in selected if row.get("split") == split]
    return float(selected[-1].get("value", selected[-1].get("accuracy", selected[-1].get("loss", "nan"))))


def comparison_results(experiment: Experiment) -> tuple[tuple[str, ...], tuple[tuple[str, ...], ...]]:
    title = experiment.title
    original_root = ROOT / "exp/ds1/results/original/data"
    if title == "MNIST Optimizer Comparison":
        folders = {
            "SGD": "dlfs1.ch06.optimizer-mnist.sgd",
            "Momentum": "dlfs1.ch06.optimizer-mnist.momentum",
            "AdaGrad": "dlfs1.ch06.optimizer-mnist.adagrad",
            "Adam": "dlfs1.ch06.optimizer-mnist.adam",
        }
        current_labels = {
            "MLP-OPT-SGD": "SGD",
            "MLP-OPT-MOMENTUM": "Momentum",
            "MLP-OPT-ADAGRAD": "AdaGrad",
            "MLP-OPT-ADAM": "Adam",
        }
        series_by_label = {label: series for series, label in current_labels.items()}
        rows = []
        for label, folder in folders.items():
            book = original_last_value(original_root / "e01" / folder, label)
            _n, value, _sd = STD_FINALS[title][series_by_label[label]]
            rows.append((label, number(book), mean_sd(title, series_by_label[label]), number(value - book)))
        return ("Condition", "Book run", "Current mean ± SD", "Difference"), tuple(rows)
    if title == "MNIST Weight Initialization":
        folders = {
            "Normal σ=0.01": ("dlfs1.ch06.init-compare.std-001", "std=0.01"),
            "Xavier": ("dlfs1.ch06.init-compare.xavier", "Xavier"),
            "He": ("dlfs1.ch06.init-compare.he", "He"),
        }
        series = {"Normal σ=0.01": "MLP-INIT-STD001", "Xavier": "MLP-INIT-XAVIER", "He": "MLP-INIT-HE"}
        rows = []
        for label, (folder, condition) in folders.items():
            book = original_last_value(original_root / "e02" / folder, condition)
            _n, value, _sd = STD_FINALS[title][series[label]]
            rows.append((label, number(book), mean_sd(title, series[label]), number(value - book)))
        return ("Condition", "Book run", "Current mean ± SD", "Difference"), tuple(rows)
    if title in {"Weight Decay and Overfitting", "Dropout and Overfitting"}:
        if title.startswith("Weight"):
            folder = original_root / "e03/dlfs1.ch06.weight-decay.lambda-01"
            prefix = "REG-WD-01/test"
            labels = (("No regularization", "Not cached", "REG-WD-OFF/test"), ("Book-default regularization", None, prefix))
        else:
            folder = original_root / "e04/dlfs1.ch06.dropout.on-ratio-02"
            prefix = "REG-DROPOUT-ON-02/test"
            labels = (("No dropout", "Not cached", "REG-DROPOUT-OFF/test"), ("Book-default dropout", None, prefix))
        rows = []
        book_value = original_last_value(folder, split="test")
        for label, unavailable, series in labels:
            _n, current, _sd = STD_FINALS[title][series]
            if unavailable:
                rows.append((label, unavailable, mean_sd(title, series, percent=True), "—"))
            else:
                rows.append((label, number(book_value, percent=True), mean_sd(title, series, percent=True), f"{(current-book_value)*100:+.2f} pp"))
        return ("Condition", "Book run", "Current mean ± SD", "Difference"), tuple(rows)
    if title == "Batch Normalization and Weight Scale":
        rows = []
        for index, row in enumerate(experiment.result_rows, 1):
            book_off = original_last_value(original_root / f"e05/dlfs1.ch06.batchnorm.scale-{index:02d}.bn-off")
            book_on = original_last_value(original_root / f"e05/dlfs1.ch06.batchnorm.scale-{index:02d}.bn-on")
            rows.append(
                (
                    row[0],
                    number(book_off, percent=True),
                    mean_sd(title, f"BN-SCALE-{index:02d}-OFF", percent=True),
                    number(book_on, percent=True),
                    mean_sd(title, f"BN-SCALE-{index:02d}-ON", percent=True),
                )
            )
        return ("Initial σ", "Book no BN", "Current no BN", "Book BN", "Current BN"), tuple(rows)
    if title in {"Simple Convolutional Network", "Deep Convolutional Network"}:
        deep = title.startswith("Deep")
        folder = original_root / ("e07/dlfs1.ch08.deep-convnet" if deep else "e06/dlfs1.ch07.simple-convnet")
        book = original_last_value(folder, split="test-full")
        current_series = "DeepCNN/test" if deep else "SimpleCNN/test"
        _n, current, _sd = STD_FINALS[title][current_series]
        return (
            ("Metric", "Book run", "Current mean ± SD", "Difference"),
            (("Final full-test accuracy", number(book, percent=True), mean_sd(title, current_series, percent=True), f"{(current-book)*100:+.2f} pp"),),
        )
    if title == "Spatial Layout Sensitivity":
        return (
            ("Condition", "Book counterpart", "Current test accuracy"),
            tuple(
                (label, "No direct experiment", mean_sd(title, series, percent=True))
                for label, series in (
                    ("MLP / original", "NN-MATCHED/test"),
                    ("MLP / permuted", "NN-MATCHED-PERMUTED/test"),
                    ("CNN / original", "CNN-SIMPLE-SPATIAL/test"),
                    ("CNN / permuted", "CNN-SIMPLE-SPATIAL-PERMUTED/test"),
                )
            ),
        )
    if title == "Optimizer Trajectories":
        folders = {"SGD": "sgd", "Momentum": "momentum", "AdaGrad": "adagrad", "Adam": "adam"}
        current = {row[0]: row for row in experiment.result_rows}
        rows = []
        for label, suffix in folders.items():
            trajectory = read_csv(original_root / f"e09/dlfs1.ch06.optimizer-path.{suffix}/trajectory.csv")[-1]
            rows.append((label, number(trajectory["x"]), number(trajectory["y"]), current[label][1], current[label][2]))
        return ("Optimizer", "Book x @30", "Book y @30", "Current x @60", "Current y @60"), tuple(rows)
    if title == "Activation Distribution Observation":
        archive = original_root / "e10/dlfs1.ch06.activation.sigmoid-std-1/activations.npz"
        with np.load(archive) as arrays:
            book_peaks = tuple(str(int(np.histogram(arrays[f"layer_{layer}"], bins=30, range=(0, 1))[0].max())) for layer in range(1, 6))
        rows = [("Book original", "Sigmoid / σ=1.0", *book_peaks)]
        for row in experiment.result_rows:
            rows.append(("Current", row[0], *row[1:]))
        return (
            ("Source", "Condition", "Layer 1", "Layer 2", "Layer 3", "Layer 4", "Layer 5"),
            tuple(rows),
        )
    if title == "First-Layer Convolution Filters":
        checkpoint = original_root / "e06/dlfs1.ch07.simple-convnet/checkpoint.npz"
        with np.load(checkpoint) as arrays:
            weights = arrays["param__W1"]
        current = experiment.result_rows[0]
        return (
            ("Source", "Shape", "Minimum", "Maximum", "Std. dev."),
            (
                ("Book trained SimpleCNN", str(tuple(weights.shape)).replace(" ", ""), number(weights.min()), number(weights.max()), number(weights.std())),
                ("Current trained SimpleCNN", current[1], current[2], current[3], current[4]),
            ),
        )
    if title == "Toy Word2Vec: CBOW vs Skip-gram":
        book = original_last_value(ROOT / "exp/ds2/results/original/data/e01/dlfs2.ch03.toy-cbow-full-softmax")
        current = {row[0]: row for row in experiment.result_rows}
        return (
            ("Architecture", "Book run", "Current mean ± SD", "Difference"),
            (
                ("CBOW", number(book), current["CBOW"][2], "seed protocols differ"),
                ("Skip-gram", "No book run", current["Skip-gram"][2], "Project extension"),
            ),
        )
    if title == "PTB Word2Vec Objectives":
        current = {row[0]: row for row in experiment.result_rows}
        return (
            ("Architecture / objective", "Book result", "Current mean ± SD", "Comparison"),
            (
                ("CBOW / negative sampling", "≈1.49", current["CBOW / negative sampling"][2], "Comparable"),
                ("Skip-gram / negative sampling", "Not cached", current["Skip-gram / negative sampling"][2], "Book condition"),
                ("CBOW / full softmax", "No book run", current["CBOW / full softmax"][2], "Project extension"),
                ("Skip-gram / full softmax", "No book run", current["Skip-gram / full softmax"][2], "Project extension"),
            ),
        )
    if title == "Small-Corpus RNN Language Model":
        current = experiment.result_rows[0][2]
        return ("Metric", "Book-code run", "Current mean ± SD", "Difference"), (("Final train perplexity", "6.10", current, f"{float(current.split(' ± ')[0])-6.10:+.2f}"),)
    if title == "Penn Treebank LSTM Language Model":
        current = experiment.result_rows[0][2]
        return ("Metric", "Book checkpoint", "Current mean ± SD", "Difference"), (("Final test perplexity", "136.08", current, f"{float(current.split(' ± ')[0])-136.078:+.2f}"),)
    if title == "Penn Treebank LM Recipes":
        current = {row[0]: row for row in experiment.result_rows}
        improved_value = float(current["Improved LSTM"][2].split()[0].replace(",", ""))
        return (
            ("Recipe", "Book result", "Current result", "Comparison"),
            (
                ("Vanilla RNN", "No combined book run", current["Vanilla RNN"][2], "Project comparison"),
                ("LSTM", "No combined book run", current["LSTM"][2], "Project comparison"),
                ("Improved LSTM", "80.83", current["Improved LSTM"][2], f"{improved_value-80.826:+.2f}"),
            ),
        )
    if title == "Addition Seq2Seq Models":
        labels = {
            "SEQA-VAN-FWD": "Vanilla / forward",
            "SEQA-VAN-REV": "Vanilla / reversed",
            "SEQA-PEEKY-FWD": "Peeky / forward",
            "SEQA-PEEKY-REV": "Peeky / reversed",
        }
        current = {label: mean_sd(title, series, percent=True) for series, label in labels.items()}
        return (
            ("Model / order", "Book result", "Current mean ± SD", "Comparison"),
            (
                ("Vanilla / forward", "12.94%", current["Vanilla / forward"], "Comparable"),
                ("Vanilla / reversed", "Not cached", current["Vanilla / reversed"], "Book condition"),
                ("Peeky / forward", "Not cached", current["Peeky / forward"], "Book condition"),
                ("Peeky / reversed", "Not cached", current["Peeky / reversed"], "Book condition"),
            ),
        )
    if title == "Date Conversion Seq2Seq Models":
        labels = {
            "SEQD-VAN-REV": "Vanilla",
            "SEQD-PEEKY-REV": "Peeky",
            "SEQD-ATTN-REV": "Attention",
        }
        current = {label: mean_sd(title, series, percent=True) for series, label in labels.items()}
        return (
            ("Model", "Book result", "Current mean ± SD", "Comparison"),
            (
                ("Vanilla", "Not cached", current["Vanilla"], "Book condition"),
                ("Peeky", "Not cached", current["Peeky"], "Book condition"),
                ("Attention", "100.00%", current["Attention"], "Comparable"),
            ),
        )
    if title == "Attention Alignment Observation":
        return (
            ("Item", "Book original", "Current implementation"),
            (
                ("Trained checkpoint", "Available", "Unavailable"),
                ("Alignment examples", "5", "0"),
                ("Attention map", "Produced", "Not produced"),
                ("Comparison status", "Reference available", "Pending completed run"),
            ),
        )
    raise KeyError(title)


def experiments(asset_dir: Path) -> tuple[Experiment, ...]:
    d1 = ROOT / "exp/ds1/results/image"
    d2 = ROOT / "exp/ds2/results/image"
    word2vec_graph = make_word2vec_graph(asset_dir / "ptb_word2vec_final_loss.png")
    activation_graph = make_activation_montage(asset_dir / "activation_montage.png")

    common_mnist = (
        ("Dataset", "MNIST, 60,000 train / 10,000 test"),
        ("Sampling", "Mini-batch sampling with replacement"),
        ("Seeds", "10 paired master seeds"),
        ("Objective", "Mean softmax cross-entropy"),
    )
    return (
        Experiment(
            "Deep Learning from Scratch 1",
            "MNIST Optimizer Comparison",
            (
                ("Model", "MLP: 784 → 100 × 4 → 10"),
                ("Activation / init.", "ReLU / He"),
                ("Optimizers", "SGD, Momentum 0.9, AdaGrad, Adam"),
                ("Learning rates", "0.01; Adam 0.001"),
            ),
            common_mnist
            + (("Batch / budget", "128 / 5 epochs, up to 2,000 updates"), ("Compute", "NumPy CPU, float64")),
            d1 / "e01_band.png",
            ("Optimizer", "Seeds", "Final loss", "Min–max"),
            band_rows(
                d1 / "e01_band.csv",
                {
                    "MLP-OPT-SGD": "SGD",
                    "MLP-OPT-MOMENTUM": "Momentum",
                    "MLP-OPT-ADAGRAD": "AdaGrad",
                    "MLP-OPT-ADAM": "Adam",
                },
            ),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "MNIST Weight Initialization",
            (
                ("Model", "MLP: 784 → 100 × 4 → 10"),
                ("Activation", "ReLU"),
                ("Initializers", "Normal σ=0.01, Xavier, He"),
                ("Optimizer", "SGD, learning rate 0.01"),
            ),
            common_mnist
            + (("Batch / budget", "128 / 5 epochs, up to 2,000 updates"), ("Compute", "NumPy CPU, float64")),
            d1 / "e02_band.png",
            ("Initializer", "Seeds", "Final loss", "Min–max"),
            band_rows(
                d1 / "e02_band.csv",
                {
                    "MLP-INIT-STD001": "Normal σ=0.01",
                    "MLP-INIT-XAVIER": "Xavier",
                    "MLP-INIT-HE": "He",
                },
            ),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "Weight Decay and Overfitting",
            (
                ("Model", "MLP: 784 → 100 × 6 → 10"),
                ("Activation / init.", "ReLU / He"),
                ("Weight decay", "0.0 vs 0.1"),
                ("Optimizer", "SGD, learning rate 0.01"),
            ),
            (
                ("Dataset", "MNIST: first 300 train / full 10,000 test"),
                ("Batch", "100, sampling with replacement"),
                ("Budget", "201 epochs, up to 601 updates"),
                ("Evaluation", "Train and test every 3 updates"),
                ("Seeds / compute", "10 paired seeds / NumPy CPU float64"),
            ),
            d1 / "e03_band.png",
            ("Condition", "Final train acc.", "Final test acc.", "Test min–max"),
            pivot_train_test(d1 / "e03_band.csv", (("REG-WD-OFF", "No weight decay"), ("REG-WD-01", "Weight decay 0.1"))),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "Dropout and Overfitting",
            (
                ("Model", "MLP: 784 → 100 × 6 → 10"),
                ("Activation / init.", "ReLU / He"),
                ("Dropout", "0.0 vs 0.2"),
                ("Optimizer", "SGD, learning rate 0.01"),
            ),
            (
                ("Dataset", "MNIST: first 300 train / full 10,000 test"),
                ("Batch", "100, sampling with replacement"),
                ("Budget", "301 epochs, up to 903 updates"),
                ("Evaluation", "Train and test at first update of each epoch"),
                ("Seeds / compute", "10 paired seeds / NumPy CPU float64"),
            ),
            d1 / "e04_band.png",
            ("Condition", "Final train acc.", "Final test acc.", "Test min–max"),
            pivot_train_test(d1 / "e04_band.csv", (("REG-DROPOUT-OFF", "No dropout"), ("REG-DROPOUT-ON-02", "Dropout 0.2"))),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "Batch Normalization and Weight Scale",
            (
                ("Model", "MLP: 784 → 100 × 5 → 10"),
                ("Activation", "ReLU"),
                ("Batch normalization", "Enabled vs disabled"),
                ("Initial weight σ", "16 log-spaced values: 1.0 to 0.0001"),
                ("Optimizer", "SGD, learning rate 0.01"),
            ),
            (
                ("Dataset", "MNIST: first 1,000 train samples"),
                ("Batch", "100, sampling with replacement"),
                ("Budget", "20 epochs"),
                ("Evaluation", "Training accuracy every 10 updates"),
                ("Seeds / compute", "10 paired seeds / NumPy CPU float64"),
            ),
            d1 / "e05_band.png",
            ("Initial σ", "Without BN", "With BN"),
            batchnorm_rows(d1 / "e05_band.csv"),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "Simple Convolutional Network",
            (
                ("Model", "Conv 1→30, 5×5 → pool → FC 100 → 10"),
                ("Initialization", "Normal σ=0.01"),
                ("Optimizer", "Adam, learning rate 0.001"),
                ("Objective", "Mean softmax cross-entropy"),
            ),
            (
                ("Dataset", "MNIST images, 60,000 train / 10,000 test"),
                ("Batch / budget", "100 / 20 epochs"),
                ("Evaluation", "1,000-sample train/test each epoch; full test at end"),
                ("Seeds", "10"),
                ("Compute", "CuPy GPU, float64"),
            ),
            d1 / "e06_e07_band.png",
            ("Model", "Seeds", "Final test acc.", "Min–max"),
            band_rows(d1 / "e06_e07_band.csv", {"SimpleCNN/test": "SimpleCNN"}, percent=True),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "Deep Convolutional Network",
            (
                ("Model", "6 convolution layers; channels 16 / 32 / 64"),
                ("Classifier", "FC 50 → 10"),
                ("Regularization", "Dropout 0.5"),
                ("Initialization", "He"),
                ("Optimizer", "Adam, learning rate 0.001"),
            ),
            (
                ("Dataset", "MNIST images, 60,000 train / 10,000 test"),
                ("Batch / budget", "100 / 20 epochs"),
                ("Evaluation", "1,000-sample train/test each epoch; full test at end"),
                ("Seeds", "10"),
                ("Compute", "CuPy GPU, float64"),
            ),
            d1 / "e06_e07_band.png",
            ("Model", "Seeds", "Final test acc.", "Min–max"),
            band_rows(d1 / "e06_e07_band.csv", {"DeepCNN/test": "DeepCNN"}, percent=True),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "Spatial Layout Sensitivity",
            (
                ("Models", "Parameter-matched MLP vs SimpleCNN"),
                ("MLP", "784 → 489 → 100 → 10, ReLU, He"),
                ("CNN", "Conv 1→30, 5×5 → pool → FC 100 → 10"),
                ("Input conditions", "Original vs fixed pixel permutation"),
                ("Optimizer", "Adam, learning rate 0.001"),
            ),
            (
                ("Dataset", "MNIST full train/test"),
                ("Batch / budget", "100 / 2 epochs"),
                ("Permutation seed", "20260808"),
                ("Evaluation", "Every 20 updates and full test each epoch"),
                ("Seeds / compute", "10 paired seeds / CuPy GPU float64"),
            ),
            d1 / "e08_band.png",
            ("Condition", "Final train acc.", "Final test acc.", "Test min–max"),
            pivot_train_test(
                d1 / "e08_band.csv",
                (
                    ("NN-MATCHED", "MLP / original"),
                    ("NN-MATCHED-PERMUTED", "MLP / permuted"),
                    ("CNN-SIMPLE-SPATIAL", "CNN / original"),
                    ("CNN-SIMPLE-SPATIAL-PERMUTED", "CNN / permuted"),
                ),
            ),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "Optimizer Trajectories",
            (
                ("Objective", "f(x,y) = x²/20 + y²"),
                ("Initial point", "(-7, 2)"),
                ("Optimizers", "SGD, Momentum, AdaGrad, Adam"),
                ("Learning rates", "0.95, 0.1, 1.5, 0.3"),
                ("Momentum", "0.9"),
            ),
            (
                ("Dataset", "Analytic quadratic; no sampled dataset"),
                ("Budget", "60 updates"),
                ("Seeds", "1"),
                ("Numerics", "NumPy CPU, float64"),
                ("Recorded output", "Parameter trajectory at every update"),
            ),
            d1 / "e09_band.png",
            ("Optimizer", "Final x", "Final y", "Final objective"),
            trajectory_rows(d1 / "e09_band.csv"),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "Activation Distribution Observation",
            (
                ("Model", "Five-layer MLP observation network"),
                ("Width / bias", "100 units per layer / 0.0"),
                ("Activations", "Sigmoid, tanh, ReLU"),
                ("Initializers", "σ=0.01, Xavier, He, σ=1.0"),
                ("Model seed", "40403"),
            ),
            (
                ("Input", "1,000 × 100 synthetic normal samples"),
                ("Input seed", "40402"),
                ("Training", "None; forward observation only"),
                ("Seeds / compute", "1 / NumPy CPU float64"),
                ("Recorded output", "Per-layer activation histograms"),
            ),
            activation_graph,
            ("Condition", "Layer 1", "Layer 2", "Layer 3", "Layer 4", "Layer 5"),
            activation_rows(d1 / "e10_band.csv"),
        ),
        Experiment(
            "Deep Learning from Scratch 1",
            "First-Layer Convolution Filters",
            (
                ("Observed models", "SimpleCNN under three input/training conditions"),
                ("Parameter", "First convolution weight tensor"),
                ("Kernel bank", "30 filters, 1 input channel, 5×5 kernel"),
                ("Visualization", "Shared grayscale weight scale"),
            ),
            (
                ("Dataset context", "MNIST original and fixed pixel permutation"),
                ("Checkpoint", "Final checkpoint"),
                ("Observation seed", "Master seed 1"),
                ("Compared conditions", "Full training, spatial, spatial-permuted"),
                ("Output", "Filter mosaics and weight statistics"),
            ),
            d1 / "e11_band.png",
            ("Condition", "Weight shape", "Minimum", "Maximum", "Std. dev."),
            filter_rows(d1 / "e11_band.csv"),
        ),
        Experiment(
            "Deep Learning from Scratch 2",
            "Toy Word2Vec: CBOW vs Skip-gram",
            (
                ("Models", "CBOW and Skip-gram"),
                ("Embedding size", "5"),
                ("Context window", "1"),
                ("Objective", "Full softmax"),
                ("Optimizer", "Adam, learning rate 0.001"),
            ),
            (
                ("Corpus", '"You say goodbye and I say hello."'),
                ("Batch / budget", "3 / 1,000 epochs"),
                ("Loss recording", "Mean every 20 updates"),
                ("Seeds", "10 paired seeds"),
                ("Compute", "CuPy GPU, float32, deterministic"),
            ),
            d2 / "e01_band.png",
            ("Architecture", "Seeds", "Final loss ± SD", "Training time"),
            summary_rows(
                d2 / "e01_summary.csv",
                {"W2V-TOY-CBOW-FULL": "CBOW", "W2V-TOY-SKIPGRAM-FULL": "Skip-gram"},
                "final_loss",
            ),
        ),
        Experiment(
            "Deep Learning from Scratch 2",
            "PTB Word2Vec Objectives",
            (
                ("Models", "CBOW and Skip-gram"),
                ("Embedding size", "100"),
                ("Context window", "5"),
                ("Objectives", "Full softmax vs negative sampling"),
                ("Negative samples", "5; unigram power 0.75"),
            ),
            (
                ("Dataset", "Penn Treebank training corpus"),
                ("Batch / budget", "100 / 10 epochs"),
                ("Optimizer", "Adam, learning rate 0.001"),
                ("Seeds", "10 paired seeds"),
                ("Compute", "CuPy GPU, float32, deterministic"),
            ),
            word2vec_graph,
            ("Architecture / objective", "Seeds", "Final loss ± SD", "Training time"),
            summary_rows(
                d2 / "e02_summary.csv",
                {
                    "W2V-PTB-CBOW-NS": "CBOW / negative sampling",
                    "W2V-PTB-SKIPGRAM-NS": "Skip-gram / negative sampling",
                    "W2V-PTB-CBOW-FULL": "CBOW / full softmax",
                    "W2V-PTB-SKIPGRAM-FULL": "Skip-gram / full softmax",
                },
                "final_loss",
            ),
        ),
        Experiment(
            "Deep Learning from Scratch 2",
            "Small-Corpus RNN Language Model",
            (
                ("Model", "Vanilla RNN language model"),
                ("Word vector / hidden", "100 / 100"),
                ("Optimizer", "SGD, learning rate 0.1"),
                ("Gradient clipping", "None"),
                ("Objective", "Temporal softmax cross-entropy"),
            ),
            (
                ("Dataset", "First 1,000 Penn Treebank tokens"),
                ("Batch / time steps", "10 / 5"),
                ("Budget", "100 epochs"),
                ("Seeds", "10"),
                ("Compute", "CuPy GPU, float32, deterministic"),
            ),
            d2 / "e03_band.png",
            ("Model", "Seeds", "Final train perplexity ± SD", "Training time"),
            summary_rows(d2 / "e03_summary.csv", {"LM-SMALL-RNN": "Vanilla RNN"}, "final_train_perplexity"),
        ),
        Experiment(
            "Deep Learning from Scratch 2",
            "Penn Treebank LSTM Language Model",
            (
                ("Model", "LSTM language model"),
                ("Word vector / hidden", "100 / 100"),
                ("Optimizer", "SGD, learning rate 20.0"),
                ("Gradient clipping", "0.25"),
                ("Objective", "Temporal softmax cross-entropy"),
            ),
            (
                ("Dataset", "Penn Treebank"),
                ("Batch / time steps", "20 / 35"),
                ("Budget", "4 epochs"),
                ("Evaluation", "Terminal test perplexity"),
                ("Seeds / compute", "10 / CuPy GPU float32 deterministic"),
            ),
            d2 / "e04_band.png",
            ("Model", "Seeds", "Final test perplexity ± SD", "Training time"),
            summary_rows(d2 / "e04_summary.csv", {"LM-LSTM": "LSTM"}, "final_test_perplexity"),
        ),
        Experiment(
            "Deep Learning from Scratch 2",
            "Penn Treebank LM Recipes",
            (
                ("Models", "Vanilla RNN, LSTM, improved LSTM"),
                ("Word vector / hidden", "650 / 650"),
                ("Improved model", "Dropout 0.5"),
                ("Optimizer", "SGD, initial learning rate 20.0"),
                ("Gradient clipping", "0.25"),
            ),
            (
                ("Dataset", "Penn Treebank"),
                ("Batch / time steps", "20 / 35"),
                ("Budget", "40 epochs"),
                ("Schedule", "Validation-based LR decay ÷4"),
                ("Evaluation", "Validation each epoch; test at end"),
            ),
            d2 / "e05_band.png",
            ("Recipe", "Seeds", "Final test perplexity ± SD", "Training time"),
            summary_rows(
                d2 / "e05_summary.csv",
                {
                    "LM-RNN-RECIPE": "Vanilla RNN",
                    "LM-LSTM-RECIPE": "LSTM",
                    "LM-BETTER-RECIPE": "Improved LSTM",
                },
                "final_test_perplexity",
            ),
        ),
        Experiment(
            "Deep Learning from Scratch 2",
            "Addition Seq2Seq Models",
            (
                ("Models", "Vanilla and Peeky sequence-to-sequence"),
                ("Word vector / hidden", "16 / 128"),
                ("Input order", "Forward vs reversed"),
                ("Optimizer", "Adam, learning rate 0.001"),
                ("Gradient clipping", "5.0"),
            ),
            (
                ("Dataset", "Synthetic addition sequences"),
                ("Split", "Legacy split, seed 1984"),
                ("Batch / budget", "128 / 25 epochs"),
                ("Evaluation", "Full-test exact match every epoch"),
                ("Seeds / compute", "10 paired / CuPy GPU float32"),
            ),
            d2 / "e06_band.png",
            ("Model / order", "Completed seeds", "Final accuracy ± SD", "Training time"),
            summary_rows(
                d2 / "e06_summary.csv",
                {
                    "SEQA-VAN-FWD": "Vanilla / forward",
                    "SEQA-VAN-REV": "Vanilla / reversed",
                    "SEQA-PEEKY-FWD": "Peeky / forward",
                    "SEQA-PEEKY-REV": "Peeky / reversed",
                },
                "final_test_accuracy",
            ),
        ),
        Experiment(
            "Deep Learning from Scratch 2",
            "Date Conversion Seq2Seq Models",
            (
                ("Models", "Vanilla, Peeky, Attention sequence-to-sequence"),
                ("Word vector / hidden", "16 / 256"),
                ("Input order", "Reversed"),
                ("Optimizer", "Adam, learning rate 0.001"),
                ("Gradient clipping", "5.0"),
            ),
            (
                ("Dataset", "Date-format conversion sequences"),
                ("Split", "Legacy split, seed 1984"),
                ("Batch / budget", "128 / 10 epochs"),
                ("Evaluation", "Full-test exact match every epoch"),
                ("Seeds / compute", "10 planned / CuPy GPU float32"),
            ),
            d2 / "e07_band.png",
            ("Model", "Completed seeds", "Final accuracy ± SD", "Training time"),
            summary_rows(
                d2 / "e07_summary.csv",
                {
                    "SEQD-VAN-REV": "Vanilla",
                    "SEQD-PEEKY-REV": "Peeky",
                    "SEQD-ATTN-REV": "Attention",
                },
                "final_test_accuracy",
            ),
        ),
        Experiment(
            "Deep Learning from Scratch 2",
            "Attention Alignment Observation",
            (
                ("Model", "Attention sequence-to-sequence"),
                ("Word vector / hidden", "16 / 256"),
                ("Input order", "Reversed"),
                ("Decoding", "Greedy"),
                ("Observation", "Decoder-to-encoder attention matrix"),
            ),
            (
                ("Dataset", "Date-format conversion test split"),
                ("Selection seed", "1984"),
                ("Requested samples", "5"),
                ("Compute", "NumPy CPU, float32, deterministic"),
                ("Dependency", "Completed trained attention checkpoint"),
            ),
            d2 / "e08_band.png",
            ("Item", "Result"),
            (
                ("Completed observation runs", "0"),
                ("Alignment matrices produced", "0"),
                ("Current graph status", "No completed runs"),
                ("Required next input", "Trained attention checkpoint"),
            ),
        ),
    )


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        relationship_id = slide_id.rId
        prs.part.drop_rel(relationship_id)
        slide_ids.remove(slide_id)


def placeholder(slide, kind: PP_PLACEHOLDER):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == kind:
            return shape
    raise ValueError(f"placeholder not found: {kind}")


def set_table_style(table) -> None:
    table_properties = table._tbl.tblPr
    for child in list(table_properties):
        if child.tag.endswith("tableStyleId"):
            table_properties.remove(child)
    table_properties.append(
        parse_xml(f"<a:tableStyleId {nsdecls('a')}>{TABLE_STYLE}</a:tableStyleId>")
    )


def fill_table(table, headers: tuple[str, ...], rows: tuple[tuple[str, ...], ...]) -> None:
    set_table_style(table)
    for column, value in enumerate(headers):
        table.cell(0, column).text = value
    for row_index, row in enumerate(rows, 1):
        for column, value in enumerate(row):
            table.cell(row_index, column).text = str(value)
    for row in table.rows:
        for cell in row.cells:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.025)
            cell.margin_bottom = Inches(0.025)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER


def add_settings_slide(prs: Presentation, experiment: Experiment) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = f"{experiment.title} — Setup"
    comparison, fixed_rows, difference = separated_setup(experiment.title)
    comparison_frame = slide.shapes.add_table(
        2,
        3,
        Inches(0.72),
        Inches(1.18),
        Inches(11.9),
        Inches(1.12),
    )
    fill_table(
        comparison_frame.table,
        ("Comparison variable", "Compared levels", "Target outcome"),
        (comparison,),
    )
    comparison_frame.table.columns[0].width = Inches(2.65)
    comparison_frame.table.columns[1].width = (
        Inches(6.35)
    )
    comparison_frame.table.columns[2].width = (
        comparison_frame.width
        - comparison_frame.table.columns[0].width
        - comparison_frame.table.columns[1].width
    )
    fixed_frame = slide.shapes.add_table(
        len(fixed_rows) + 1,
        2,
        Inches(0.72),
        Inches(2.48),
        Inches(11.9),
        Inches(2.95),
    )
    fill_table(
        fixed_frame.table,
        ("Fixed component", "Held-constant protocol"),
        fixed_rows,
    )
    fixed_frame.table.columns[0].width = Inches(2.7)
    fixed_frame.table.columns[1].width = fixed_frame.width - fixed_frame.table.columns[0].width
    caption = slide.shapes.add_textbox(
        Inches(0.82),
        Inches(5.66),
        Inches(11.7),
        Inches(0.72),
    )
    caption.text_frame.word_wrap = True
    caption.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    caption.text_frame.margin_left = Inches(0.04)
    caption.text_frame.margin_right = Inches(0.04)
    caption.text_frame.margin_top = Inches(0.02)
    caption.text_frame.margin_bottom = Inches(0.02)
    caption.text = difference


def add_picture_contained(slide, image_path: Path) -> None:
    left, top, width, height = Inches(0.72), Inches(1.08), Inches(11.9), Inches(5.55)
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    box_ratio = width / height
    if image_ratio > box_ratio:
        picture_width = width
        picture_height = int(width / image_ratio)
        picture_left = left
        picture_top = top + (height - picture_height) // 2
    else:
        picture_height = height
        picture_width = int(height * image_ratio)
        picture_top = top
        picture_left = left + (width - picture_width) // 2
    slide.shapes.add_picture(str(image_path), picture_left, picture_top, picture_width, picture_height)


def add_graph_slide(prs: Presentation, experiment: Experiment, graph: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = f"{experiment.title} — Results"
    add_picture_contained(slide, graph)


def add_result_table_slide(prs: Presentation, experiment: Experiment) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = f"{experiment.title} — Summary"
    headers, rows = comparison_results(experiment)
    left, top, width, height = Inches(0.72), Inches(1.22), Inches(11.9), Inches(5.15)
    frame = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        left,
        top,
        width,
        height,
    )
    fill_table(frame.table, headers, rows)


def add_section_slide(prs: Presentation, section: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = section
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            shape.text = subtitle
            break


def build() -> Path:
    apply_plot_theme()
    prs = Presentation(TEMPLATE)
    remove_all_slides(prs)
    with TemporaryDirectory(prefix="ds-presentation-") as directory:
        asset_dir = Path(directory)
        items = experiments(asset_dir)
        STD_FINALS.clear()
        standard_deviation_graphs = make_sd_current_graphs(asset_dir)
        references = book_graphs(asset_dir)
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        placeholder(cover, PP_PLACEHOLDER.CENTER_TITLE).text = "Book Reproduction Comparison"
        placeholder(cover, PP_PLACEHOLDER.SUBTITLE).text = (
            "Deep Learning from Scratch 1 & 2\n"
            "Book original vs current implementation: setups, result graphs, and summary tables"
        )
        current_section = None
        for index, experiment in enumerate(items):
            if experiment.section != current_section:
                current_section = experiment.section
                subtitle = (
                    "Optimization, initialization, regularization, convolution, and representation observations"
                    if current_section.endswith("1")
                    else "Word representation, language modeling, sequence-to-sequence, and attention"
                )
                add_section_slide(prs, current_section, subtitle)
            add_settings_slide(prs, experiment)
            comparison_graph = make_comparison_graph(
                references[experiment.title],
                standard_deviation_graphs.get(experiment.title, experiment.graph),
                asset_dir / f"comparison_{index:02d}.png",
            )
            add_graph_slide(prs, experiment, comparison_graph)
            add_result_table_slide(prs, experiment)
        prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
